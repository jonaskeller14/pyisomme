from pyisomme.plotting import *

from pptx.dml.color import RGBColor
import os


class Page:
    name: str

    def __init__(self, report):
        self.report = report

    def set(self):
        pass

    def construct(self, presentation):
        pass

    def __repr__(self):
        return f"Page({self.name})"


class Page_Cover(Page):
    name: str = "Cover"
    title: str = None
    subtitle: str = None

    def __init__(self, report):
        super().__init__(report)
        self.title = report.title
        self.subtitle = f'{report.name}\n{" | ".join([isomme.test_number for isomme in report.isomme_list])}'

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title
        slide.placeholders[1].text = self.subtitle


class Page_Result_Table(Page):
    name: str
    title: str
    table_content: dict

    def __init__(self, report):
        super().__init__(report)
        self.table = {}

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        shapes = slide.shapes

        rows = len(self.table_content) + 1
        cols = 1 + len(self.report.isomme_list)

        shape = shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table

        tbl = shape._element.graphic.graphicData.tbl
        tbl[0][-1].text = '{5940675A-B579-460E-94D1-54222C63F5DA}'

        # write column headings
        for idx, isomme in enumerate(self.report.isomme_list):
            table.cell(0, idx+1).text = isomme.test_number
            table.cell(0, idx+1).fill.solid()
            table.cell(0, idx+1).fill.fore_color.rgb = RGBColor(0xFB, 0x8F, 0x00)

        # write body cells
        for row, (row_name, row_values) in zip(range(1, rows), self.table_content.items()):
            table.cell(row, 0).text = row_name
            for idx, value in enumerate(row_values):
                table.cell(row, idx+1).text = value


class Page_Plot_nxn(Page):
    name: str
    title: str
    codes: dict
    nrows: int = 1
    ncols: int = 1
    sharey: bool = None

    def __init__(self, report):
        super().__init__(report)

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        fig = Plot_Line(self.report.isomme_list, self.codes, nrows=self.nrows, ncols=self.ncols, sharey=self.sharey, limits=self.report.limits, figsize=(figsize_x, figsize_y)).fig
        fig.savefig("tmp.png", transparent=True)

        slide.shapes.add_picture("tmp.png", left=left, top=top, height=height)

        os.remove("tmp.png")


class Page_OLC(Page_Plot_nxn):
    name: str = "OLC"
    title: str = "Occupant Load Criterion (OLC)"
    codes: dict
    nrows: int = 1
    ncols: int = 2
    sharey: bool = None

    def __init__(self, report):
        super().__init__(report)
        self.codes = {isomme: ["14BPIL??????VEXA", "14BPIL0OLC??VEXA"] for isomme in self.report.isomme_list}

