from __future__ import annotations

from pyisomme import Channel, Isomme
from pyisomme.limits import limit_list_sort
from pyisomme.plotting import Plot_Line
from pyisomme.report.criterion import Criterion

from pptx.dml.color import RGBColor
import os
import matplotlib.pyplot as plt
from matplotlib.colors import to_rgb
import numpy as np
import io


class Page:
    name: str

    def __init__(self, report):
        self.report = report

    def set(self):
        pass

    def construct(self, presentation):
        pass

    def __repr__(self):
        return f"Page({self.name})"


class Page_Cover(Page):
    name: str = "Cover"
    title: str = None
    subtitle: str = None

    def __init__(self, report):
        super().__init__(report)
        self.title = report.title
        self.subtitle = f'{report.name}\n{" | ".join([isomme.test_number for isomme in report.isomme_list])}'

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title
        slide.placeholders[1].text = self.subtitle


class Page_Criterion_Values_Table(Page):
    # TODO: matplotlib table
    name: str
    title: str
    criteria: dict[Isomme, list[Criterion]]

    def __init__(self, report):
        super().__init__(report)
        self.criteria = {}

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        shapes = slide.shapes

        rows = len(list(self.criteria.values())[0]) + 1
        cols = 1 + len(self.report.isomme_list)

        shape = shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table

        tbl = shape._element.graphic.graphicData.tbl
        tbl[0][-1].text = '{5940675A-B579-460E-94D1-54222C63F5DA}'

        # write column headings
        for idx, isomme in enumerate(self.report.isomme_list):
            table.cell(0, idx+1).text = isomme.test_number
            table.cell(0, idx+1).fill.solid()
            table.cell(0, idx+1).fill.fore_color.rgb = RGBColor(200, 200, 200)

        # write body cells
        for idx_isomme, isomme in enumerate(self.report.isomme_list):
            for idx_criterion, criterion in enumerate(self.criteria[isomme]):
                if idx_isomme == 0:
                    table.cell(idx_criterion+1, 0).text = f"{criterion.name} [{criterion.channel.unit if criterion.channel is not None else np.nan}]"
                table.cell(idx_criterion+1, idx_isomme+1).text = f"{criterion.value:.4g}"

                if criterion.color is not None:
                    table.cell(idx_criterion+1, idx_isomme+1).fill.solid()
                    table.cell(idx_criterion+1, idx_isomme+1).fill.fore_color.rgb = RGBColor(*[int(255 * x) for x in to_rgb(criterion.color)])


class Page_Criterion_Rating_Table(Page):
    # TODO matplotlib table
    name: str
    title: str
    criteria: dict[Isomme, list[Criterion]]

    def __init__(self, report):
        super().__init__(report)
        self.criteria = {}

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        shapes = slide.shapes

        rows = len(list(self.criteria.values())[0]) + 1
        cols = 1 + len(self.report.isomme_list)

        shape = shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table

        tbl = shape._element.graphic.graphicData.tbl
        tbl[0][-1].text = '{5940675A-B579-460E-94D1-54222C63F5DA}'

        # write column headings
        for idx, isomme in enumerate(self.report.isomme_list):
            table.cell(0, idx+1).text = isomme.test_number
            table.cell(0, idx+1).fill.solid()
            table.cell(0, idx+1).fill.fore_color.rgb = RGBColor(200, 200, 200)

        # write body cells
        for idx_isomme, isomme in enumerate(self.report.isomme_list):
            for idx_criterion, criterion in enumerate(self.criteria[isomme]):
                if idx_isomme == 0:
                    table.cell(idx_criterion+1, 0).text = f"{criterion.name}"
                table.cell(idx_criterion+1, idx_isomme+1).text = f"{criterion.rating:.2g}"

                if criterion.color is not None:
                    table.cell(idx_criterion+1, idx_isomme+1).fill.solid()
                    table.cell(idx_criterion+1, idx_isomme+1).fill.fore_color.rgb = RGBColor(*[int(255 * x) for x in to_rgb(criterion.color)])


class Page_Criterion_Values_Chart(Page):
    name: str
    title: str
    criteria: dict[Isomme, list[Criterion]]

    def __init__(self, report):
        super().__init__(report)
        self.criteria = {}

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        fig, ax = plt.subplots(figsize=(figsize_x, figsize_y), layout="constrained")

        x_labels = [c.name for c in list(self.criteria.values())[0]]
        col_factors = np.ones(len(x_labels))
        line_values = np.array([[abs(c.value) for c in self.criteria[isomme]] for isomme in self.report.isomme_list])
        # TODO: units --> conversion of channel and limit values, take units from first variant?

        # Calculate Column Factor
        for idx, criterion in enumerate(self.criteria[self.report.isomme_list[0]]):
            col_factors[idx] = np.nanmax([np.nanmax([abs(l.func(0)) for l in criterion.limits.limit_list]), np.nanmax(np.abs(line_values[:, idx]))]) * 1.1

        # Plot Bars
        for idx_col, criterion in enumerate(self.criteria[self.report.isomme_list[0]]):
            limits = limit_list_sort(criterion.limits.limit_list, sym=True)
            limit_values = [abs(limit.func(0)) for limit in limits]
            for idx, (limit, limit_value) in enumerate(zip(limits, limit_values)):
                if idx == 0:
                    ax.bar(x=criterion.name,
                           height=limit_value / col_factors[idx_col],
                           color=limit.color,
                           alpha=0.5)
                elif idx < len(limits) - 1:
                    if (limit.lower and limit.func(0) >= 0) or (limit.upper and limit.func(0) < 0):
                        ax.bar(x=criterion.name,
                               bottom=limit_value / col_factors[idx_col],
                               height=(limit_values[idx+1] - limit_value) / col_factors[idx_col],
                               color=limit.color,
                               alpha=0.5)
                    if (limit.upper and limit.func(0) >= 0) or (limit.lower and limit.func(0) < 0):
                        ax.bar(x=criterion.name,
                               bottom=limit_values[idx-1] / col_factors[idx_col],
                               height=(limit_value - limit_values[idx-1]) / col_factors[idx_col],
                               color=limit.color,
                               alpha=0.5)
                elif idx == len(limits) - 1:
                    ax.bar(x=criterion.name,
                           bottom=limit_value / col_factors[idx_col],
                           height=1-limit_value / col_factors[idx_col],
                           color=limit.color,
                           alpha=0.5)

        # Plot Lines
        for idx, isomme in enumerate(self.report.isomme_list):
            ax.plot(x_labels, line_values[idx, :] / col_factors, marker="o", label=isomme.test_number, linewidth=3, markersize=8)

        plt.xticks(rotation=30, ha='right')
        ax.get_yaxis().set_visible(False)
        ax.legend()

        image_steam = io.BytesIO()
        fig.savefig(image_steam, transparent=True)
        slide.shapes.add_picture(image_steam, left=left, top=top, height=height)


class Page_Plot_nxn(Page):
    name: str
    title: str
    channels: dict[Isomme, list[list[Channel | str]]]
    nrows: int = 1
    ncols: int = 1
    sharey: bool = False

    def __init__(self, report):
        super().__init__(report)

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        fig = Plot_Line(self.channels, nrows=self.nrows, ncols=self.ncols, sharey=self.sharey, limits=self.report.limits, figsize=(figsize_x, figsize_y)).fig

        image_steam = io.BytesIO()
        fig.savefig(image_steam, transparent=True)
        slide.shapes.add_picture(image_steam, left=left, top=top, height=height)


class Page_OLC(Page_Plot_nxn):
    name: str = "OLC"
    title: str = "Occupant Load Criterion (OLC)"
    channels: dict
    nrows: int = 1
    ncols: int = 1

    def __init__(self, report):
        super().__init__(report)
        self.channels = {isomme: [[isomme.get_channel("14BPIL??????VEXA", "10SEATLERE??VEXA"),
                                   isomme.get_channel("14BPIL0OLC??VEXA", "10SEAT0OLC??VEXA")]] for isomme in self.report.isomme_list}
